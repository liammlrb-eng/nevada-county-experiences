"""
Generate demo_pitch.pptx — a focused 15-slide deck for the county presentation.

Output is .pptx, which:
  - Opens directly in Google Slides (Drive → upload → "Open with Google Slides")
  - Opens in PowerPoint, Keynote, LibreOffice Impress
  - Stays editable; all text/colors live in the source script

Run:   python generate_slides.py
Out:   demo_pitch.pptx

ORGANIZATION
============
Each slide is a self-contained builder function (slide_title, slide_strategic_insight,
slide_what_this_is, etc.) that takes a slide object and adds shapes/text to it.
The narrative order is defined in SLIDE_ORDER at the bottom — reordering the deck
is just reordering that list.

Page footers (X / N) and TOTAL are computed automatically from SLIDE_ORDER, so
nothing breaks when slides move.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

# ── Brand palette — modern ─────────────────────────────────────────────
# Deep ink-blue as the primary dark; vivid amber as the warm accent; teal
# as the cool accent. Replaces the earth-tone palette to feel less
# brochure-y and more contemporary on a projector. BROWN/GOLD names kept
# for backward compat — they now resolve to ink/amber.
BROWN  = RGBColor(0x10, 0x24, 0x44)          # INK — deep slate-blue for headers (was warm brown)
GOLD   = RGBColor(0xF2, 0xA9, 0x3E)          # AMBER — vivid warm accent (was muted gold)
GOLD_LIGHT = RGBColor(0xFF, 0xD1, 0x8C)      # warm amber tint
RUST   = RGBColor(0xF8, 0x71, 0x71)          # CORAL — bold callout accent (rarely used)
FOREST = RGBColor(0x14, 0xB8, 0xA6)          # TEAL — cool accent for variety
SLATE  = RGBColor(0x33, 0x41, 0x55)          # modern dark text
FOG    = RGBColor(0xF1, 0xF5, 0xF9)          # cool light wash (was warm cream)
DARK   = RGBColor(0x0F, 0x17, 0x2A)          # near-black ink
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
RULE   = RGBColor(0xCB, 0xD5, 0xE1)          # cool light divider line

# Global font-size multiplier. Applied at every font.size = Pt(...) site
# below so a single number controls the deck-wide legibility floor.
# 1.30 = ~30% larger across the board.
FONT_SCALE = 1.30

# Show or hide the '↳ scenario' italic sub-lines under every bullet and
# table row. Off by default — the headline bullet alone is enough; the
# scenarios crowd the layout and read as filler at projector distance.
SHOW_SCENARIOS = False

# 16:9 widescreen
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def add_text(slide, text, left, top, width, height, *,
             size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT,
             font='Calibri', anchor=MSO_ANCHOR.TOP, italic=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size * FONT_SCALE)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def add_rect(slide, left, top, width, height, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.5)
    shp.shadow.inherit = False
    return shp


def add_gold_rule(slide, left, top, width):
    """Title underline. Name kept for backward compat but the actual color
    is now deep ink (matching the heading) so rules look like a heading
    underline rather than a gold ribbon. Amber is still used elsewhere
    for callouts / badges where it pops against white."""
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(2.5))
    shp.fill.solid()
    shp.fill.fore_color.rgb = BROWN
    shp.line.fill.background()


def slide_header(slide, title, subtitle=None):
    """Standard top-of-slide header: title, gold rule, optional subtitle."""
    add_text(slide, title, Inches(0.6), Inches(0.4), Inches(12.1), Inches(0.7),
             size=32, bold=True, color=BROWN, font='Calibri')
    add_gold_rule(slide, Inches(0.6), Inches(1.18), Inches(12.1))
    if subtitle:
        add_text(slide, subtitle, Inches(0.6), Inches(1.32), Inches(12.1), Inches(0.5),
                 size=14, color=SLATE, italic=True)


def page_footer(slide, page_num, total):
    """Footer with page number + brand line."""
    add_text(slide, 'Western Nevada County Experience',
             Inches(0.6), Inches(7.1), Inches(8), Inches(0.3),
             size=9, color=SLATE)
    add_text(slide, f'{page_num} / {total}',
             Inches(11.5), Inches(7.1), Inches(1.2), Inches(0.3),
             size=9, color=SLATE, align=PP_ALIGN.RIGHT)
    add_gold_rule(slide, Inches(0.6), Inches(7.06), Inches(12.1))


def add_bullets(slide, items, left, top, width, height, *,
                size=14, color=DARK, line_spacing=1.4, bullet_color=GOLD,
                scenario_size=None, scenario_color=None):
    """
    Each item is either a string OR a (text, scenario) tuple.
    Strings render as a single bulleted line.
    Tuples render the item, then an italic gold sub-line "↳ scenario"
    showing how the feature is used in practice.

    scenario_size  defaults to size - 2
    scenario_color defaults to BROWN (the brand brown)
    """
    if scenario_size is None:
        scenario_size = max(size - 2, 8)
    if scenario_color is None:
        scenario_color = BROWN
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)

    para_idx = 0
    for item in items:
        # Allow either plain string or (text, scenario) tuple
        if isinstance(item, tuple):
            text, scenario = item
        else:
            text, scenario = item, None

        # Bullet line
        p = tf.paragraphs[0] if para_idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        b = p.add_run()
        b.text = '• '
        b.font.size = Pt((size + 2) * FONT_SCALE)
        b.font.bold = True
        b.font.color.rgb = bullet_color
        b.font.name = 'Calibri'
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size * FONT_SCALE)
        r.font.color.rgb = color
        r.font.name = 'Calibri'
        p.line_spacing = line_spacing
        if para_idx > 0:
            p.space_before = Pt(6 if scenario is None else 4)
        para_idx += 1

        # Optional scenario sub-line: italic, indented, smaller
        if scenario and SHOW_SCENARIOS:
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.LEFT
            sr = p2.add_run()
            sr.text = '       ↳  ' + scenario
            sr.font.size = Pt(scenario_size * FONT_SCALE)
            sr.font.italic = True
            sr.font.color.rgb = scenario_color
            sr.font.name = 'Calibri'
            p2.line_spacing = 1.1
            p2.space_before = Pt(0)
            para_idx += 1
    return box


def add_stat_card(slide, left, top, width, height, big_text, label):
    """Featured number + label, like a metric tile."""
    add_rect(slide, left, top, width, height, FOG, RULE)
    # Big number
    add_text(slide, big_text, left, top + Inches(0.2), width, Inches(0.9),
             size=44, bold=True, color=GOLD, align=PP_ALIGN.CENTER, font='Calibri')
    # Label
    add_text(slide, label, left, top + Inches(1.05), width, Inches(0.5),
             size=11, color=SLATE, align=PP_ALIGN.CENTER, font='Calibri')


def two_col_table(slide, headers, rows, top, *, left_w=4.5, right_w=8.0,
                  header_size=11, body_size=10, row_height=0.5,
                  scenario_size=None):
    """
    Simple 2-column table without using actual table shape (more layout control).
    Each row is either a 2-tuple (left, right) or a 3-tuple (left, right, scenario).
    A scenario renders as a second italic line in the right column.
    """
    if scenario_size is None:
        scenario_size = max(body_size - 2, 8)
    left  = Inches(0.6)
    col2  = left + Inches(left_w)
    # Header row
    add_rect(slide, left, top, Inches(left_w + right_w), Inches(0.42), FOG)
    add_text(slide, headers[0], left + Inches(0.1), top + Inches(0.05),
             Inches(left_w - 0.15), Inches(0.32),
             size=header_size, bold=True, color=BROWN)
    add_text(slide, headers[1], col2 + Inches(0.1), top + Inches(0.05),
             Inches(right_w - 0.15), Inches(0.32),
             size=header_size, bold=True, color=BROWN)
    add_gold_rule(slide, left, top + Inches(0.42), Inches(left_w + right_w))
    y = top + Inches(0.46)
    for i, row in enumerate(rows):
        if len(row) == 3:
            a, b, scenario = row
        else:
            a, b = row
            scenario = None
        # Honor the global hide-scenarios flag — even 3-tuples render as
        # just the headline pair when SHOW_SCENARIOS is False.
        if not SHOW_SCENARIOS:
            scenario = None

        if i % 2 == 0:
            add_rect(slide, left, y, Inches(left_w + right_w),
                     Inches(row_height), FOG)

        # Left column — vertically centered when row is tall
        add_text(slide, a, left + Inches(0.1), y + Inches(0.05),
                 Inches(left_w - 0.15), Inches(row_height - 0.05),
                 size=body_size, color=DARK,
                 anchor=MSO_ANCHOR.MIDDLE if scenario else MSO_ANCHOR.TOP)

        if scenario:
            # Right column — explanation on top, scenario in italic below
            half = (row_height - 0.05) / 2
            add_text(slide, b, col2 + Inches(0.1), y + Inches(0.04),
                     Inches(right_w - 0.15), Inches(half),
                     size=body_size, color=DARK)
            add_text(slide, '↳  ' + scenario,
                     col2 + Inches(0.1), y + Inches(half + 0.04),
                     Inches(right_w - 0.15), Inches(half),
                     size=scenario_size, color=BROWN, italic=True)
        else:
            add_text(slide, b, col2 + Inches(0.1), y + Inches(0.05),
                     Inches(right_w - 0.15), Inches(row_height - 0.05),
                     size=body_size, color=DARK)
        y += Inches(row_height)


# ═══════════════════════════ SLIDE BUILDERS ═══════════════════════════════
# Each function takes a slide object and adds content to it (no page_footer,
# no add_slide call — those happen in the main loop). Order is defined by
# SLIDE_ORDER below.

def slide_title(s):
    # Brown band on left
    add_rect(s, 0, 0, Inches(4.5), SLIDE_H, BROWN)
    # Gold accent
    add_rect(s, Inches(4.5), 0, Inches(0.08), SLIDE_H, GOLD)
    add_text(s, 'Western\nNevada\nCounty',
             Inches(0.6), Inches(2.2), Inches(3.7), Inches(2.5),
             size=44, bold=True, color=GOLD_LIGHT, font='Calibri')
    add_text(s, 'Trip Planner',
             Inches(0.6), Inches(4.4), Inches(3.7), Inches(0.7),
             size=22, color=WHITE, italic=True, font='Calibri')
    add_text(s, 'Visitor planning platform — county pitch deck',
             Inches(4.9), Inches(3.0), Inches(7.8), Inches(0.6),
             size=20, color=BROWN, bold=True)
    add_gold_rule(s, Inches(4.9), Inches(3.6), Inches(7.5))
    add_text(s,
        'A purpose-built tourism platform combining 161 curated experiences '
        'with 1,200+ live local events. Privacy-first, opt-in, no tracking.',
        Inches(4.9), Inches(3.85), Inches(7.8), Inches(2.5),
        size=14, color=SLATE)


def slide_strategic_insight(s):
    """The lodging-revenue framing — leads the deck so everything else has context."""
    slide_header(s, 'The strategic insight',
        'What sets this apart from a generic tourism site.')

    add_text(s,
        'Most tourism sites optimize for visits to the site. The right metric is '
        'lodging bookings driven by the site. Those are different optimization targets.',
        Inches(0.6), Inches(1.85), Inches(12.1), Inches(1.0),
        size=15, color=DARK)

    add_bullets(s, [
        ('Day-by-day itinerary view structurally surfaces the "where do I sleep?" question',
         'Visitor sees Day 2 has no hotel and books one before leaving the page.'),
        ('"Find Lodging" CTA at the top of every itinerary — booking always one click away',
         'No hunting for a button — booking decision is visible on every itinerary view.'),
        ('Empty lodging slots create a visual gap that drives the booking decision',
         'That blank Day 2 stay tile is psychologically jarring; the visitor fills it.'),
        ('New items default to the latest day — no day-1 pile-up; trip naturally fills out',
         'Adding stops naturally extends the trip rather than cramming Day 1.'),
        ('Print view goes long-form per day — visitors who print, follow it',
         'Couple prints the itinerary; the same Day 1 / Day 2 structure follows on paper.'),
        ('Smart suggestions are tag + geography aware — clusters keep visitors close to lodging',
         'After Holbrooke Hotel goes in, suggestions cluster around Mill Street.'),
    ], Inches(0.6), Inches(2.85), Inches(12.1), Inches(3.5),
       size=11, line_spacing=1.2, scenario_size=9)

    # Closing quote
    add_rect(s, Inches(0.6), Inches(6.2), Inches(12.1), Inches(0.7),
             FOG, GOLD)
    add_text(s,
        '"Convince a visitor to add one more night and you double or triple local revenue."',
        Inches(0.7), Inches(6.25), Inches(11.9), Inches(0.6),
        size=13, italic=True, color=BROWN, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def slide_what_this_is(s):
    slide_header(s, 'What this is',
        'A tourism platform for Western Nevada County — Gold Country foothills.')

    # Body intro
    add_text(s,
        'A purpose-built tourism platform for Western Nevada County — '
        'centered on Nevada City and Grass Valley, with full coverage of Penn Valley, '
        'North San Juan, Rough & Ready, Washington, Chicago Park, Smartsville, and '
        'adjacent Colfax. The platform combines curated places with live local '
        'events, so visitors discover what to do AND when it\'s actually happening.',
        Inches(0.6), Inches(1.95), Inches(12.1), Inches(1.4),
        size=13, color=DARK)

    # Stat cards
    add_stat_card(s, Inches(0.6),  Inches(3.5), Inches(2.85), Inches(1.6), '161',  'Curated experiences\nacross 11 communities')
    add_stat_card(s, Inches(3.65), Inches(3.5), Inches(2.85), Inches(1.6), '1,200+', 'Live events from\n10+ sources')
    add_stat_card(s, Inches(6.7),  Inches(3.5), Inches(2.85), Inches(1.6), '9',    'Themed vibes\nfor discovery')
    add_stat_card(s, Inches(9.75), Inches(3.5), Inches(2.85), Inches(1.6), '$5–$20', 'Typical\nmonthly cost')

    # Below: scope footer
    add_text(s,
        'Geographic scope: Western Nevada County (Gold Country foothills) + adjacent Colfax.\n'
        'Truckee and Sierra-side communities are out of scope for this build.',
        Inches(0.6), Inches(5.4), Inches(12.1), Inches(0.7),
        size=11, color=SLATE, italic=True)


def slide_discovery(s):
    slide_header(s, 'How visitors discover',
        'Search, browse, or filter by date — multiple paths to "I found it".')

    two_col_table(s, ['Discovery path', 'What it does'], [
        ('Universal keyword search',
         'One box searches every venue and event — names, descriptions and tags — with built-in synonyms, so a search for "metalsmith" also finds a studio listed only as "blacksmithing"',
         'A visitor types "blacksmithing", clicks Find, and lands on The Curious Forge at once.'),
        ('🌿 In Season Now strip',
         'A scrolling banner under the hero showing what\'s happening RIGHT NOW — wildflowers peaking, swim season open, festivals starting — with "Peak now / Just started / Ends soon" badges',
         'October visitor lands and sees "Aspens are on fire right now" — books the trip on the spot.'),
        ('Themed vibes (9 cards)',
         'Historic · Arts · Hands-On · Foodie · Active · Relaxed · Wellness · Family · Festivals',
         'No-plan visitor taps "Foodie" — relevant cards instantly filter in.'),
        ('🏠 Lodging tile — direct route to stays',
         'A 10th tile alongside the vibes; bypasses the vibe-funnel for "where do I sleep" intent',
         'First-time visitor sees the house-shaped Lodging tile, jumps straight to 50+ stays.'),
        ('Vibe-level pills (12+ across all vibes)',
         'Foodie → Wineries / Restaurants / Markets · Active → Hiking / Pickleball / Day-Pass Gyms',
         'Active-vibe visitor taps "🏋️ Day-Pass Gyms" — finds 8 places to keep their routine going.'),
        ('Category dropdown + sub-pills',
         'Lodging → Hotels / B&Bs / Glamping / Campgrounds / RV  (and more)',
         'Family of four picks Lodging → B&Bs to find character stays over chains.'),
        ('Activity tags on every card',
         'Hiking · Biking · Swimming · Fishing · Boating · Running visible at a glance',
         'Mountain biker scrolls and immediately spots cards tagged "Biking".'),
        ('Date filter + auto-season',
         'All Future · Today · Weekend · This Week — auto-hides summer-only items in winter',
         'December visitor never sees "Summer Concert Series" — it filters itself out.'),
        ('Area filter',
         'County-wide events show for any local area; city events filter normally',
         'Penn Valley resident filters out Nevada City clutter for a local night out.'),
        ('Tag-aware Smart Suggestions',
         '"Near Your Stops" panel in the itinerary — geographic + tag scoring',
         'After Empire Mine goes in, the panel surfaces Holbrooke Hotel 0.4 mi away.'),
    ], Inches(1.55), left_w=3.4, right_w=8.7, body_size=9, row_height=0.50,
       scenario_size=7.5)


def slide_help_me_plan(s):
    """The trip generator — six questions, one click, a real day-by-day plan."""
    slide_header(s, '✨ Help Me Plan — one-click trip generator',
        'Six quick questions → a complete day-by-day itinerary made from real local venues.')

    add_text(s,
        'Removes the blank-page problem for first-time visitors. The planner scores '
        'every venue in the catalog against the visitor\'s answers and assembles a '
        'real trip — not generic suggestions, not AI hallucination, just rule-based '
        'matching against curated local data.',
        Inches(0.6), Inches(1.85), Inches(12.1), Inches(1.1),
        size=13, color=DARK)

    # Left: what the visitor answers
    add_text(s, 'What the visitor answers',
             Inches(0.6), Inches(3.05), Inches(6), Inches(0.4),
             size=14, bold=True, color=GOLD)
    add_bullets(s, [
        'Trip length (1–5 days)',
        'Pace — Relaxed / Balanced / Packed',
        'Who\'s coming — Family / Couple / Friends / Solo / Multi-gen',
        'Activity level — Easy / Moderate / Adventurous',
        'Indoor vs outdoor lean',
        'Interests (Foodie · Historic · Active · etc.)',
        'Optional: must-haves (live music, hike, wine tasting…)',
        'Optional: trip start date + chosen hotel',
    ], Inches(0.6), Inches(3.5), Inches(6), Inches(3.0),
       size=11, line_spacing=1.3)

    # Right: what the visitor gets
    add_text(s, 'What the visitor gets',
             Inches(7.0), Inches(3.05), Inches(5.7), Inches(0.4),
             size=14, bold=True, color=GOLD)
    add_bullets(s, [
        ('A complete day-by-day itinerary, ready to edit',
         'No blank Day 1 — every day is pre-populated with real stops.'),
        ('Stops biased toward the chosen hotel',
         'Pick Holbrooke → most stops cluster near Grass Valley.'),
        ('Matching events woven into the right day',
         'Concert on visit Friday lands in Day 1 automatically.'),
        ('Group-aware picks (kids, romance, adventure)',
         'Family inputs → kid-friendly venues; couple → romantic.'),
        ('Editable: change anything, add anything, remove anything',
         'Visitor tweaks the plan in the same My Itinerary they\'d use anyway.'),
    ], Inches(7.0), Inches(3.5), Inches(5.7), Inches(3.0),
       size=11, line_spacing=1.2, scenario_size=9)

    # Closing pitch line
    add_rect(s, Inches(0.6), Inches(6.6), Inches(12.1), Inches(0.5),
             FOG, GOLD)
    add_text(s,
        'Rule-based, not AI-generated — no hallucinated venues, no API costs. The output is always a real plan from real local data.',
        Inches(0.7), Inches(6.65), Inches(11.9), Inches(0.4),
        size=11, italic=True, color=BROWN, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def slide_itinerary(s):
    slide_header(s, 'Day-by-day itinerary builder',
        'Visitors think in nights, not stops — the planner reflects that.')

    # Left: feature list
    add_bullets(s, [
        ('Each day has its own card with a numbered header and calendar date',
         'Visitor sees "Day 1: Sat May 17" instead of an undifferentiated stop list.'),
        ('Per-day "Tonight\'s Stay" lodging slot — empty slots prompt for booking',
         'Empty Day 2 stay tile reminds the visitor they haven\'t booked Sunday yet.'),
        ('Move stops between days (dropdown) + drag-to-reorder within a day',
         'Spouse on phone moves Friday\'s hike to Saturday in two taps.'),
        ('Add Day → returns visitor to browse mode for the new day',
         'Tap "Add Day," land on cards to pick Day 3 activities — flow continues.'),
        ('Find Lodging button at top of itinerary jumps directly to lodging filter',
         'Six stops in but no hotel — one tap surfaces every nearby B&B.'),
        ('Events go INTO the itinerary first; the visitor decides what to commit to',
         'Bluegrass concert tile lands in My Itinerary, not on KVMR\'s site.'),
        ('📅 Set trip dates → events auto-weave to the right day + Calendar view available',
         'Pick Jun 14–15; Saturday\'s bluegrass concert lands in Day 1 automatically.'),
        ('🗺 Map view — every stop on a Leaflet map; click pins to navigate',
         'Visitor opens the map: sees the full weekend laid out, decides Day 2 needs a lunch stop near the trail.'),
        ('🖨 Print itinerary — clean PDF-ready layout for paper handouts',
         'Older visitor prints the day-by-day plan, walks Mill Street with paper in hand.'),
        ('Opt-in save survives tab close + browser restart — no account, no install, no email',
         'Browser closes, phone reboots — itinerary still there; no signup wall.'),
        ('Share link via native phone share sheet, email, text, or copy URL',
         'Visitor texts the URL to spouse: "this is what I\'m thinking."'),
    ], Inches(0.6), Inches(1.75), Inches(7.0), Inches(5.3),
       size=10, line_spacing=1.15, scenario_size=8)

    # Right: visual mock representation
    mock_x = Inches(7.9); mock_y = Inches(1.85); mock_w = Inches(4.8); mock_h = Inches(5.0)
    add_rect(s, mock_x, mock_y, mock_w, mock_h, FOG, RULE)
    add_text(s, '2 stops · 2 days / 1 night',
             mock_x + Inches(0.2), mock_y + Inches(0.15), mock_w - Inches(0.4), Inches(0.4),
             size=11, color=SLATE, italic=True)
    # Day 1 mock
    add_rect(s, mock_x + Inches(0.25), mock_y + Inches(0.6), mock_w - Inches(0.5), Inches(2.0),
             WHITE, GOLD)
    add_text(s, 'DAY 1 of 2   Sat, Jun 14',
             mock_x + Inches(0.4), mock_y + Inches(0.7), mock_w - Inches(0.6), Inches(0.3),
             size=11, bold=True, color=BROWN)
    add_text(s, '💤 TONIGHT\'S STAY\nHolbrooke Hotel · Grass Valley',
             mock_x + Inches(0.4), mock_y + Inches(1.1), mock_w - Inches(0.6), Inches(0.7),
             size=10, color=DARK)
    add_text(s, '1.  Empire Mine State Historic Park\n2.  Lola Restaurant',
             mock_x + Inches(0.4), mock_y + Inches(1.85), mock_w - Inches(0.6), Inches(0.7),
             size=10, color=DARK)
    # Day 2 mock
    add_rect(s, mock_x + Inches(0.25), mock_y + Inches(2.75), mock_w - Inches(0.5), Inches(1.6),
             WHITE, GOLD)
    add_text(s, 'DAY 2 of 2   Sun, Jun 15',
             mock_x + Inches(0.4), mock_y + Inches(2.85), mock_w - Inches(0.6), Inches(0.3),
             size=11, bold=True, color=BROWN)
    add_text(s, '💤  No lodging picked for tonight  →  [Find Lodging]',
             mock_x + Inches(0.4), mock_y + Inches(3.2), mock_w - Inches(0.6), Inches(0.5),
             size=9, italic=True, color=SLATE)
    add_text(s, '📅  Father\'s Day Bluegrass — Sat 7:00 PM',
             mock_x + Inches(0.4), mock_y + Inches(3.7), mock_w - Inches(0.6), Inches(0.5),
             size=10, color=GOLD)


def slide_demo_flow(s):
    slide_header(s, 'Let me show you',
        'A live demo: visitor planning a romantic November weekend in Western Nevada County.')

    acts = [
        ('Act 0 (optional) · 60 sec', '✨ Help Me Plan opener — answer six questions, full draft itinerary appears. Use for blank-page audiences; skip for power-users who want to see manual planning.'),
        ('Act 1 · 2 min',  'Set the scene — couple from Sacramento, click "This Weekend" + Festivals vibe'),
        ('Act 2 · 90 sec', 'Layer in Relaxed vibe + Lodging category + B&B sub-pill — real Victorian B&Bs surface'),
        ('Act 3 · 2 min',  'Build the weekend — Cornish Christmas + Holbrooke Hotel + Empire Mine + Lola Restaurant. Open My Itinerary — Smart Suggestions appear with "0.4 mi away" labels. Map view + Print + Share link.'),
        ('Act 4 · 90 sec', 'Behind the scenes — admin Events Queue, Approve All, AI Categorize button, Event Sources tab'),
        ('Act 5 · 60 sec', 'Close the value props — time saved, local visibility, privacy, no vendor lock-in, AI integration, low maintenance'),
    ]
    two_col_table(s, ['Act', 'What happens'], acts,
                  Inches(1.65), left_w=2.3, right_w=9.8,
                  body_size=10, row_height=0.75)


def slide_personas(s):
    slide_header(s, 'Who this serves',
        'Nine visitor personas — each searches differently and converts on different content.')

    personas = [
        ('1. Romantic Weekender',     'Couples 35-65, Bay Area / Sac, B&Bs + dining + wineries'),
        ('2. Festival Pilgrim',       'Anchored to a specific event — Cornish Christmas, Wild & Scenic, Bluegrass'),
        ('3. Trail & Lake Seeker',    'Hiking / biking / swimming / paddling — needs distance & difficulty'),
        ('4. Multi-Gen Family',       'Parents with kids 5-15 — "will my 8-year-old be bored?"'),
        ('5. Arts & Heritage Traveler', 'Substance over photos; reads notes; values authenticity'),
        ('6. Wine Country Foodie',    'Napa-alternative — provenance, chef bios, tasting flights'),
        ('7. Wellness Refugee',       'Float, massage, yoga, gardens — filtering for what\'s NOT there'),
        ('8. Local "What\'s Up"',     'Residents — pure events feed, "what\'s on this weekend"'),
        ('9. Maker Traveler',         'Travels for the workshop — a destination-class workshop venue can anchor a regional identity'),
    ]
    add_bullets(s, [f'{name}    —    {desc}' for name, desc in personas],
                Inches(0.6), Inches(1.85), Inches(12.1), Inches(4.7),
                size=12, line_spacing=1.4)

    # Strategic call-out
    add_rect(s, Inches(0.6), Inches(6.0), Inches(12.1), Inches(0.85),
             FOG, GOLD)
    add_text(s,
        'Recommended marketing focus (over-serve four):  Festival Pilgrim  ·  Romantic Weekender  ·  Foodie  ·  Maker Traveler',
        Inches(0.7), Inches(6.05), Inches(11.9), Inches(0.75),
        size=12, bold=True, color=BROWN, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def slide_ux_decisions(s):
    slide_header(s, 'Why it\'s easy to use',
        'Friction points where most tourism sites lose visitors — and what we did instead.')

    ux_pairs = [
        ('Day-by-day with lodging anchor',
         'A flat list of stops hides the night question; structured days surface it',
         'Visitor sees Day 2 has no stay tile, books before leaving the page.'),
        ('"Find Lodging" CTA + auto-clear vibes',
         'One-click jump to lodging that doesn\'t fight against active vibes',
         'Foodie-vibe visitor taps Find Lodging; vibes clear so all stays show.'),
        ('Add Day returns to browse',
         'No leaving visitor staring at empty days; flow continues naturally',
         'Day 2 isn\'t an empty container — visitor lands somewhere productive.'),
        ('Events → itinerary first',
         'No hijack to KVMR — visitor compares before clicking through',
         'Three concerts compared side-by-side before committing to one.'),
        ('Vibe sub-pills',
         'Foodie → Restaurants/Markets/Wineries — granular without leaving the vibe',
         'Foodie → Wineries narrows without switching mental contexts.'),
        ('Calendar dates on day headers',
         '"DAY 1 · Sat, May 17" instead of abstract numbers',
         '"Sat, May 17" grounds the trip in real-world weekend planning.'),
        ('Opt-in itinerary save',
         'Survives tab close + restart — no account, no email, explicit consent',
         'Banner appears only after first add — visitor consents knowingly.'),
        ('Auto-season filtering',
         'November visit hides summer-only experiences automatically',
         'December visitor never sees "summer concert series" as a dead lead.'),
    ]
    two_col_table(s, ['Decision', 'Why it matters'], ux_pairs,
                  Inches(1.65), left_w=4.0, right_w=8.1, body_size=9,
                  row_height=0.6, scenario_size=8)


def slide_privacy(s):
    slide_header(s, 'Privacy posture',
        'No tracking. No cookies. Visitors stay anonymous; the chamber gets credit for it.')

    add_bullets(s, [
        ('No analytics, no cookies, no third-party trackers by default',
         'Site loads instantly — no cookie modal between visitor and content.'),
        ('Itinerary save is opt-in only — explicit consent banner on first add',
         'Banner appears only after first card added; visitor knows what they\'re consenting to.'),
        ('"Forget on this device" link revokes consent and wipes saved data anytime',
         'Trip ends — visitor wipes the itinerary in two taps, takes nothing back home.'),
        ('Share-link URL works without any consent — purely a URL fragment',
         'Two friends compare itineraries via shared URLs, neither registered for anything.'),
        ('No account / sign-up / email required to plan a full trip',
         '30 seconds from landing → planning, no email-collection wall.'),
        ('Chamber sees zero personal information about visitors',
         'Aggregate event queue tells the chamber what\'s hot; visitor identity stays local.'),
        ('No GDPR / CCPA banners required (because nothing is collected)',
         'EU visitor on holiday gets the site without a banner ambush.'),
        ('No privacy review, no consent-platform fees, no compliance overhead',
         'Annual compliance budget for the platform: $0 — nothing to review.'),
        ('Opt-in measurement available if needed (Plausible — no cookies, no PII)',
         'If the chamber wants aggregate metrics later, flip a config flag — no rebuild.'),
    ], Inches(0.6), Inches(1.7), Inches(12.1), Inches(4.7),
       size=11, line_spacing=1.18, scenario_size=9)

    # Closing quote
    add_rect(s, Inches(0.6), Inches(6.0), Inches(12.1), Inches(0.85),
             FOG, GOLD)
    add_text(s,
        '"Visitors stay anonymous; the chamber gets credit for respecting them."',
        Inches(0.7), Inches(6.05), Inches(11.9), Inches(0.75),
        size=14, italic=True, color=BROWN, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def slide_admin(s):
    slide_header(s, 'Behind the scenes — low-admin operations',
        '~30 minutes per week of chamber-staff time keeps the site live and current.')

    add_bullets(s, [
        'One-click event-source updates from 11+ sources — KVMR, Eventbrite, NC/GV Chambers, Go Nevada, The Union, Nevada Theatre, Curious Forge, Crazy Horse Saloon, Golden Era Lounge, Wolf Craft Collective (and growing)',
        'Approve / dismiss events in a queue; bulk-approve is one click',
        'AI Categorize button (Claude Haiku) — refines area, venue, tags, quality (~$0.20 per full run)',
        'Publish and edit chamber-curated experiences via the public planner — no separate admin tool',
        'Inline-editable experience table — anyone who can edit a spreadsheet can maintain it',
        'Tag taxonomy editor — add, rename, delete tags without code',
        'Event-source URL management — add a new source by pasting a URL',
        'Pruning past events on every scrape, plus a manual button between scrapes',
        'Public RSS feed at /feed.rss — partners republish without integration work',
        'Public "Suggest a venue or event" form — submissions queue for chamber review',
    ], Inches(0.6), Inches(1.75), Inches(12.1), Inches(5.3),
       size=14, line_spacing=1.35)


def slide_ai(s):
    slide_header(s, 'AI-powered tag refinement',
        'Claude Haiku adds intelligence where keyword rules stop short — at trivial cost.')

    # Left: what AI fixes
    add_text(s, 'What it fixes', Inches(0.6), Inches(1.85), Inches(6), Inches(0.45),
             size=14, bold=True, color=GOLD)
    add_bullets(s, [
        ('Nearly every KVMR event tagged generic "Nevada County" → AI infers actual community',
         '"KVMR Storytelling Night" gets "Nevada City" instead of vague "Nevada County".'),
        ('Empty location field on most KVMR events → AI extracts venue name from description',
         '"Center for the Arts" surfaces as the venue from a description-only event.'),
        ('"Center for the Arts" → Grass Valley; "Miners Foundry" → Nevada City',
         'Visitor filtering by Grass Valley gets concerts without manual tagging.'),
        ('Cluttered descriptions → clean one-line summaries',
         'Long press-release blob shrinks to a clean teaser the visitor will actually read.'),
        ('Truckee / Sierra-side events flagged "low quality" → auto-hidden',
         'Tahoe events tagged "low quality" auto-hide on the public site, no manual cleanup.'),
        ('Future source refreshes auto-categorize new events with the same logic',
         'Tuesday\'s new event gets correct area + venue + tags without human touch.'),
    ], Inches(0.6), Inches(2.3), Inches(6), Inches(4.2),
       size=10, line_spacing=1.15, scenario_size=8)

    # Right: cost table
    add_text(s, 'Cost', Inches(7.2), Inches(1.85), Inches(5.5), Inches(0.45),
             size=14, bold=True, color=GOLD)
    two_col_table(s, ['Run pattern', 'Monthly'], [
        ('First-time bulk (full catalog)',  '~$1 once'),
        ('Daily refresh + categorize',      '$0.50 – $1.50'),
        ('Weekly refresh + categorize',     '$0.20 – $0.60'),
        ('Realistic chamber operation',     '~$0.50 – $1.50'),
    ], Inches(2.3), left_w=3.0, right_w=2.4, body_size=10, row_height=0.45)

    # Footer note
    add_rect(s, Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.5),
             FOG, GOLD)
    add_text(s,
        'Hard $5/month spending cap settable in Anthropic console — cannot be exceeded.',
        Inches(0.7), Inches(6.55), Inches(11.9), Inches(0.4),
        size=11, italic=True, color=BROWN, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def slide_cost(s):
    slide_header(s, 'What it costs',
        'Three operating scenarios for the first year of operation.')

    # 4 stat cards: Lean, Typical, Premium, Comparison
    add_stat_card(s, Inches(0.6),  Inches(1.85), Inches(2.85), Inches(1.7),
                  '$0–$50',  'LEAN\nFree hosting + AI off')
    add_stat_card(s, Inches(3.65), Inches(1.85), Inches(2.85), Inches(1.7),
                  '$80–$150', 'TYPICAL\n$5/mo VPS + weekly AI')
    add_stat_card(s, Inches(6.7),  Inches(1.85), Inches(2.85), Inches(1.7),
                  '$300–$500', 'PREMIUM\nManaged + daily AI')
    add_stat_card(s, Inches(9.75), Inches(1.85), Inches(2.85), Inches(1.7),
                  '$5K–$30K',  'COMMERCIAL ALT.\nCrowdRiff / Simpleview')

    # Bullet list of what's NOT a cost
    add_text(s, 'What\'s NOT a cost',
             Inches(0.6), Inches(3.85), Inches(6), Inches(0.45),
             size=14, bold=True, color=GOLD)
    add_bullets(s, [
        'Per-listing fees — every chamber member free',
        'Per-event fees — event sources cost zero',
        'Per-visitor fees — traffic 10× has no impact',
        'Vendor termination fees — owned outright',
        'Privacy compliance overhead — nothing collected to comply about',
        'License fees — fully open-source stack',
    ], Inches(0.6), Inches(4.3), Inches(12.1), Inches(2.6),
       size=12, line_spacing=1.4)


def slide_next_steps(s):
    """
    Retitled per chat-Claude feedback: "What the chamber needs to do next"
    sounded like assigning homework. New framing is collaborative.
    """
    slide_header(s, 'Where the group can take this',
        'Findings from the build — collaborative work that would unlock the next version.')

    # Three columns of action items grouped by category
    col_y = Inches(1.85); col_h = Inches(4.6)

    # ── Column 1: Data partnerships ──
    add_text(s, 'DATA PARTNERSHIPS',
             Inches(0.6), col_y, Inches(4.0), Inches(0.4),
             size=12, bold=True, color=GOLD)
    add_bullets(s, [
        'Ask NCAC for read access to their Trumba '
        'master — its 540-event, 16-month window is '
        'the most comprehensive feed in the county',
        'Reach out to GoNevadaCounty (gonevadacounty.com) '
        'for a public events feed — Cloudflare blocks '
        'automated access and we lose the 13 anchor festivals',
        'Coordinate with The Union for a public events '
        'feed — articles are paywalled and the RSS often '
        'returns nothing usable',
        'Confirm whether the County itself runs a Trumba '
        'master calendar — if so, federate from it instead '
        'of NCAC alone',
    ], Inches(0.6), col_y + Inches(0.4), Inches(4.0), col_h - Inches(0.4),
       size=9, line_spacing=1.35)

    # ── Column 2: Member outreach ──
    add_text(s, 'MEMBER OUTREACH',
             Inches(4.7), col_y, Inches(4.0), Inches(0.4),
             size=12, bold=True, color=GOLD)
    add_bullets(s, [
        'Encourage venues to publish events on a '
        'platform with a public feed — Trumba, Tribe '
        'Events (WordPress), Eventbrite — not just '
        'static HTML pages',
        'Push schema.org/Event JSON-LD adoption — one '
        'snippet on each event page makes auto-importing '
        'reliable without custom code per site',
        'Promote the public "Suggest a venue or event" '
        'form to members — already wired (footer + admin '
        'queue); the lift is socialization, not engineering',
        'Identify the 5-10 venues responsible for the bulk '
        'of visitor traffic and prioritize their data '
        'reliability over long-tail coverage',
    ], Inches(4.7), col_y + Inches(0.4), Inches(4.0), col_h - Inches(0.4),
       size=9, line_spacing=1.35)

    # ── Column 3: Operations ──
    add_text(s, 'OPERATIONS',
             Inches(8.8), col_y, Inches(4.0), Inches(0.4),
             size=12, bold=True, color=GOLD)
    add_bullets(s, [
        'Schedule the scraper to run nightly via Windows '
        'Task Scheduler / cron / systemd — auto-prune of '
        'past events is built in but only fires when the '
        'scraper runs',
        'Designate a queue curator — someone clears the '
        'pending-events queue 1-2x per week (Approve / '
        'Dismiss takes ~5 minutes per session)',
        'Click 🧹 Prune Past between scrapes for ad-hoc '
        'cleanup — same logic as the auto-prune, useful '
        'when events get added or dismissed off-cycle',
        'Run AI Categorize after each refresh — fixes '
        'venue, area, tags on long-tail events ($0.50/mo)',
        'Validate "no event" alerts — when a source '
        'returns 0 events, it has either gone down or '
        'restructured (NCAC was a false negative for '
        'months before we caught it)',
        'Quarterly review of disabled sources — vendors '
        'change platforms (e.g. NCAC moved to Trumba) and '
        'previously-blocked sources become reachable',
    ], Inches(8.8), col_y + Inches(0.4), Inches(4.0), col_h - Inches(0.4),
       size=8.5, line_spacing=1.3)

    # Bottom callout: what we already learned
    add_rect(s, Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.5),
             FOG, GOLD)
    add_text(s,
        'Concrete win from this work:  NCAC\'s calendar went from "0 events" '
        '(blocked by an iframe) to 540 events — by hitting Trumba\'s JSON feed directly.',
        Inches(0.7), Inches(6.55), Inches(11.9), Inches(0.4),
        size=10, italic=True, color=BROWN, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def slide_migration(s):
    """
    Per chat-Claude: "9-phase migration plan for a prototype conversation feels
    premature." Demoted to the appendix — pull up only if IT is in the room or
    someone explicitly asks how this would actually be hosted.
    """
    # Appendix eyebrow — small caps top-right so it doesn't fight the title
    add_text(s, 'APPENDIX  ·  PULL UP IF ASKED',
             Inches(8.5), Inches(0.45), Inches(4.2), Inches(0.3),
             size=9, bold=True, color=GOLD, align=PP_ALIGN.RIGHT)
    slide_header(s, 'Migration to a county server',
        'One focused half-day for someone comfortable with Linux; two days for a learner.')

    steps = [
        ('Phase 0',  'Decisions: hosting (DigitalOcean $6/mo recommended), domain, OS, server account ownership'),
        ('Phase 1',  'Pre-flight: domain registered, DNS access, SSH keys ready, Anthropic API key generated'),
        ('Phase 2',  'Provision: $6/mo Ubuntu droplet, lockdown SSH, install Python + nginx + chromium'),
        ('Phase 3',  'Deploy: git clone, Python venv, install deps, drop API keys in the config file'),
        ('Phase 4',  'systemd service for auto-restart on crash, auto-start on reboot'),
        ('Phase 5',  'nginx reverse proxy + Let\'s Encrypt HTTPS (auto-renew via certbot)'),
        ('Phase 6',  'Schedule nightly: source refresh at 3 AM + AI categorize at 3:30 AM (refresh triggers auto-prune; skip this step and expired events pile up in the queue)'),
        ('Phase 7',  'Backups: provider snapshots ($1/mo) + daily JSON tar to off-site'),
        ('Phase 8',  'Test checklist: site loads, vibes work, events populate, AI button works, source updates run'),
        ('Phase 9',  'Hand-off to chamber: admin URL, billing alerts, cheat sheet, emergency contact'),
    ]
    two_col_table(s, ['Phase', 'What happens'], steps,
                  Inches(1.85), left_w=1.6, right_w=10.5,
                  body_size=10, row_height=0.4)


def slide_thank_you(s):
    # Brown band on left
    add_rect(s, 0, 0, Inches(4.5), SLIDE_H, BROWN)
    add_rect(s, Inches(4.5), 0, Inches(0.08), SLIDE_H, GOLD)

    add_text(s, 'Thank you',
             Inches(0.6), Inches(2.4), Inches(3.7), Inches(1.0),
             size=44, bold=True, color=GOLD_LIGHT, font='Calibri')
    add_text(s, 'Questions?',
             Inches(0.6), Inches(3.4), Inches(3.7), Inches(0.7),
             size=22, color=WHITE, italic=True, font='Calibri')

    add_text(s, 'The leave-behind one-liner',
             Inches(4.9), Inches(2.4), Inches(7.8), Inches(0.5),
             size=14, bold=True, color=GOLD)
    add_gold_rule(s, Inches(4.9), Inches(2.85), Inches(7.5))
    add_text(s,
        '"This isn\'t a website — it\'s a planning workspace for visitors and a '
        'living directory for the chamber, with AI-assisted curation built in. '
        'Privacy-respecting by default, owned by the county outright, '
        'maintained in minutes per week."',
        Inches(4.9), Inches(3.05), Inches(7.8), Inches(2.5),
        size=14, italic=True, color=BROWN)

    # Footer info
    add_text(s,
        'Live demo:  liammlrb-eng.github.io/nevada-county-experiences\n'
        'Repo:       github.com/liammlrb-eng/nevada-county-experiences\n'
        'Pitch deck: demo_pitch.pdf  ·  Operator guide: operator_guide.pdf',
        Inches(4.9), Inches(5.6), Inches(7.8), Inches(1.5),
        size=11, color=SLATE, font='Consolas')


# ═══════════════════════════ DECK NARRATIVE ═══════════════════════════════
# Five-act structure based on chat-Claude's reorder feedback (May 2026):
#
#   Act 1 — Why we're here     (title + strategic insight)
#   Act 2 — What I built       (what this is + discovery + itinerary + demo)
#   Act 3 — Why this works     (personas + UX + privacy)
#   Act 4 — How it runs        (admin + AI + cost)
#   Act 5 — What's next        (group next steps + migration appendix + close)
#
# Each entry: (function, show_page_footer)
SLIDE_ORDER = [
    # ── Act 1 ─────────────────────────────────────────────────────────────
    (slide_title,              False),
    (slide_strategic_insight,  True),   # Lead with the lodging-revenue framing
    # ── Act 2 ─────────────────────────────────────────────────────────────
    (slide_what_this_is,       True),   # Numbers as evidence for Act 1's argument
    (slide_discovery,          True),
    (slide_help_me_plan,       True),   # Trip generator — biggest new feature
    (slide_itinerary,          True),
    (slide_demo_flow,          True),   # Live demo — the emotional peak
    # ── Act 3 ─────────────────────────────────────────────────────────────
    (slide_personas,           True),   # After demo, personas feel like real people
    (slide_ux_decisions,       True),
    (slide_privacy,            True),
    # ── Act 4 ─────────────────────────────────────────────────────────────
    (slide_admin,              True),
    (slide_ai,                 True),
    (slide_cost,               True),   # Second big moment — $6/mo vs $5K–$30K
    # ── Act 5 ─────────────────────────────────────────────────────────────
    (slide_next_steps,         True),   # "Where the group can take this"
    (slide_thank_you,          False),
]

# Appendix slides — shown after the close. Pull up only if asked
# (chat-Claude: "9-phase migration plan for a prototype conversation feels
# premature"). Numbered with an 'A' prefix instead of counting against TOTAL.
APPENDIX_ORDER = [
    (slide_migration,          True),
]

TOTAL = len(SLIDE_ORDER)


# ═══════════════════════════ BUILD THE DECK ═══════════════════════════════
for i, (build_fn, show_footer) in enumerate(SLIDE_ORDER, 1):
    s = prs.slides.add_slide(BLANK)
    build_fn(s)
    if show_footer:
        page_footer(s, i, TOTAL)

# Appendix — shown after Thank You so they're not in the main flow
for ai, (build_fn, show_footer) in enumerate(APPENDIX_ORDER, 1):
    s = prs.slides.add_slide(BLANK)
    build_fn(s)
    if show_footer:
        page_footer(s, f'A{ai}', f'A{len(APPENDIX_ORDER)}')


# ─── Save ───────────────────────────────────────────────────────────────
import shutil
project_root = Path(__file__).resolve().parent
out = project_root / 'demo_pitch.pptx'
prs.save(str(out))
print(f'Wrote: {out}')
print(f'Slides: {TOTAL} main + {len(APPENDIX_ORDER)} appendix = {TOTAL + len(APPENDIX_ORDER)} total')
print(f'Size: {out.stat().st_size:,} bytes ({out.stat().st_size/1024:.1f} KB)')

# Mirror to pitch_resources/ so the chat-upload bundle stays in sync
mirror_dir = project_root / 'pitch_resources'
if mirror_dir.is_dir():
    mirror = mirror_dir / 'demo_pitch.pptx'
    shutil.copy2(out, mirror)
    print(f'Mirrored: {mirror}')
